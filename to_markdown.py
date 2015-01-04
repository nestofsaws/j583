#!/usr/bin/env/python

"""
A very rough script for converting a pptx file to markdown format for the
purposes of using a presentation system such as reveal.js

Requires python-pptx
"""

from __future__ import unicode_literals
import sys

from pptx import Presentation
from pptx.shapes.picture import Picture


SLIDE_SEPARATOR = '@@'

IMAGE_PREFIX = 'images/image'
IMAGE_SUFFIX = '.png'


def get_runs(presentation):
    runs = []
    img_index = 1
    for slide_index, slide in enumerate(presentation.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                if isinstance(shape, Picture):
                    runs.append((slide_index, True, False, '{0}{1}{2}'.format(IMAGE_PREFIX, img_index, IMAGE_SUFFIX)))
                    img_index += 1
                    continue
                else:
                    continue
            for par_index, paragraph in enumerate(shape.text_frame.paragraphs):
                for run in paragraph.runs:
                    header = True if par_index == shape_index == 0 else False
                    if header:
                        last = len(runs) - 1
                        try:
                            existing = runs[last]
                        except IndexError:
                            existing = None
                        if existing is not None and existing[2]:
                            existing_text = existing[3]
                            new = (existing[0], existing[1], True, existing_text + run.text)
                            runs[last] = new
                        else:
                            runs.append((slide_index, False, True, run.text))
                    else:
                        runs.append((slide_index, False, False, run.text))
    return runs


def as_markdown(runs):
    for run_num, run in enumerate(runs):
        slide_num, is_image, is_header, text = run
        if is_image:
            print("![]({0})".format(text))
            print("")
            continue
        if is_header:
            if run_num > 0:
                print(SLIDE_SEPARATOR)
                print("")
            print("###{0}".format(text))
            print("")
            continue
        print(text)
        print("")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: Must provide a pptx file as an argument")
        sys.exit(1)
    prs = Presentation(sys.argv[1])
    runs = get_runs(prs)
    as_markdown(runs)
